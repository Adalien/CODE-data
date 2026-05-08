# -*- coding: utf-8 -*-
"""AI週會心得報告 - 採購單與採購對帳單"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Color palette - Deep Navy + Teal
NAVY     = RGBColor(0x1E, 0x27, 0x61)
DNAV     = RGBColor(0x12, 0x18, 0x40)
TEAL     = RGBColor(0x02, 0x80, 0x90)
SEAFOAM  = RGBColor(0x02, 0xC3, 0x9A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWH    = RGBColor(0xF4, 0xF7, 0xFA)
DKGRAY   = RGBColor(0x33, 0x33, 0x33)
MIDGRAY  = RGBColor(0x66, 0x66, 0x77)
LTBLUE   = RGBColor(0xCA, 0xDC, 0xFC)
GOLD     = RGBColor(0xF9, 0xE7, 0x95)
CORAL    = RGBColor(0xF9, 0x61, 0x67)
GRNBG    = RGBColor(0xEE, 0xF7, 0xF9)
TEALBG   = RGBColor(0x02, 0x40, 0x55)

EI = 914400
W, H = 13.33, 7.5

def e(v): return round(v * EI)
def ep(v): return round(v * 12700)

prs = Presentation()
prs.slide_width  = e(W)
prs.slide_height = e(H)

def bl(): return prs.slides.add_slide(prs.slide_layouts[6])

def rc(sl, l, t, w, h, c):
    s = sl.shapes.add_shape(1, e(l), e(t), e(w), e(h))
    s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s

def tx(sl, l, t, w, h, text, sz, c, bold=False, al=PP_ALIGN.LEFT, fn='Microsoft JhengHei'):
    bx = sl.shapes.add_textbox(e(l), e(t), e(w), e(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = c
    r.font.bold = bold; r.font.name = fn
    return bx

def hdr(sl, title, n=None):
    rc(sl, 0, 0, W, 1.0, NAVY)
    rc(sl, W-2.5, 0, 2.5, 1.0, TEAL)
    rc(sl, 0, 0, 0.18, H, TEAL)
    if n:
        tx(sl, W-2.2, 0.15, 1.8, 0.70, f'{n:02d}', 30, WHITE, bold=True, al=PP_ALIGN.CENTER)
    tx(sl, 0.4, 0.13, W-3.0, 0.74, title, 20, WHITE, bold=True)

# ─── SLIDE 1: Cover ──────────────────────────────────────────────────────────
sl = bl()
rc(sl, 0, 0, W, H, NAVY)
rc(sl, W-3.8, 0, 3.8, H, DNAV)
rc(sl, 0, 0, 0.3, H, TEAL)
rc(sl, 0.3, 3.1, W-0.3, 0.05, SEAFOAM)
tx(sl, 0.6, 1.1, W-4.5, 1.1, 'Claude AI 應用心得報告', 36, WHITE, bold=True)
tx(sl, 0.6, 2.3, W-4.5, 0.6, '採購單與採購對帳單流程優化', 20, LTBLUE)
tx(sl, 0.6, 3.3, 6.0, 0.5, '達特世', 16, SEAFOAM)
tx(sl, 0.6, 3.85, 6.0, 0.45, '廠務助理／採購  ADA', 15, LTBLUE)
tx(sl, 0.6, 4.35, 6.0, 0.45, '2026 年 4 月', 14, MIDGRAY)
rc(sl, W-3.5, 1.5, 3.1, 4.2, TEALBG)
tx(sl, W-3.3, 1.65, 2.7, 0.48, '本月主題', 13, SEAFOAM, bold=True, al=PP_ALIGN.CENTER)
for i, item in enumerate(['採購單比對', '對帳單核查', 'AI 自動分析', '實驗成果']):
    tx(sl, W-3.2, 2.25+i*0.68, 2.6, 0.55, '▶  ' + item, 14, WHITE)

# ─── SLIDE 2: Background ─────────────────────────────────────────────────────
sl = bl()
hdr(sl, '為什麼用 AI 處理採購作業？', 1)
rows = [
    ('採購單與對帳單數量龐大，人工核對耗時費力'),
    ('資料格式不一（紙本 / 系統截圖 / 手機拍照），容易出錯'),
    ('月底對帳壓力大，需快速確認每家廠商差異'),
    ('導入 Claude AI，協助自動化比對與整理，節省大量時間'),
]
icons = ['📦', '📄', '📅', '🤖']
for i, txt in enumerate(rows):
    t = 1.20 + i * 1.18
    bg = OFFWH if i % 2 == 0 else WHITE
    rc(sl, 0.3, t, W-0.5, 1.05, bg)
    rc(sl, 0.45, t+0.12, 0.65, 0.70, TEAL if i % 2 == 0 else SEAFOAM)
    tx(sl, 0.47, t+0.15, 0.63, 0.65, icons[i], 22, WHITE, al=PP_ALIGN.CENTER)
    tx(sl, 1.28, t+0.18, W-1.6, 0.68, txt, 16, DKGRAY)
rc(sl, 0.3, 1.20+4*1.18+0.08, W-0.5, 0.55, TEAL)
tx(sl, 0.5, 1.20+4*1.18+0.15, W-0.7, 0.38,
   '目標：讓採購對帳從「2～3 小時」縮短至「30 秒以內」', 15, WHITE, bold=True)

# ─── SLIDE 3: Purchase Orders ─────────────────────────────────────────────────
sl = bl()
hdr(sl, '採購單 — Claude 幫了什麼？', 2)
rc(sl, 0.3, 1.1, 6.1, 5.95, OFFWH)
rc(sl, 0.3, 1.1, 6.1, 0.45, NAVY)
tx(sl, 0.5, 1.17, 5.7, 0.35, '▌ AI 輔助採購單處理', 14, WHITE, bold=True)
left_items = [
    ('快速整理供應商報價，統一格式輸出'),
    ('自動比對品號、規格、單價是否一致'),
    ('異常品項自動標記，節省人工查核時間'),
    ('耳繩庫存、原物料庫存資料自動解析'),
]
marks = ['①', '②', '③', '④']
for i, txt in enumerate(left_items):
    t = 1.70 + i * 1.12
    rc(sl, 0.48, t, 0.50, 0.50, TEAL)
    tx(sl, 0.48, t, 0.50, 0.50, marks[i], 16, WHITE, bold=True, al=PP_ALIGN.CENTER)
    tx(sl, 1.12, t+0.03, 5.15, 0.80, txt, 15, DKGRAY)
rc(sl, 6.7, 1.1, 6.3, 5.95, GRNBG)
rc(sl, 6.7, 1.1, 6.3, 0.45, SEAFOAM)
tx(sl, 6.9, 1.17, 5.9, 0.35, '▌ 傳統 vs AI 對比', 14, WHITE, bold=True)
pairs = [
    ('傳統：人工逐行核對', 'AI：自動全覽比對'),
    ('傳統：格式不統一需手整', 'AI：統一輸出 Excel'),
    ('傳統：容易遺漏異常', 'AI：自動標記異常'),
]
for i, (lft, rgt) in enumerate(pairs):
    t = 1.70 + i * 1.65
    rc(sl, 6.85, t, 2.8, 1.45, CORAL)
    tx(sl, 6.92, t+0.35, 2.65, 0.70, lft, 13, WHITE, al=PP_ALIGN.CENTER)
    rc(sl, 9.90, t, 2.85, 1.45, SEAFOAM)
    tx(sl, 9.97, t+0.35, 2.70, 0.70, rgt, 13, WHITE, al=PP_ALIGN.CENTER)
    tx(sl, 9.60, t+0.45, 0.28, 0.55, '→', 18, TEAL, bold=True, al=PP_ALIGN.CENTER)
rc(sl, 0.3, 7.02, W-0.5, 0.40, NAVY)
tx(sl, 0.5, 7.07, W-0.7, 0.32,
   '本月已完成：喬煒興業、福綿、凱達印刷 三家廠商採購單比對', 13, WHITE, bold=True)

# ─── SLIDE 4: Reconciliation ─────────────────────────────────────────────────
sl = bl()
hdr(sl, '採購對帳單 — 實際效益', 3)
sw = (W-0.9)/4 - 0.13
for i, (step, desc, clr) in enumerate([
    ('STEP 1', '從 ERP\n匯出資料', NAVY),
    ('STEP 2', 'AI 自動\n彙整分析', TEAL),
    ('STEP 3', '比對進貨單\n與發票金額', NAVY),
    ('STEP 4', '輸出對帳\n摘要報告', SEAFOAM),
]):
    lf = 0.35 + i * (sw + 0.13)
    rc(sl, lf, 1.15, sw, 2.0, clr)
    tc = WHITE if clr != SEAFOAM else DKGRAY
    tx(sl, lf, 1.22, sw, 0.45, step, 13, tc, bold=True, al=PP_ALIGN.CENTER)
    tx(sl, lf+0.1, 1.72, sw-0.2, 1.0, desc, 16, tc, bold=True, al=PP_ALIGN.CENTER)
    if i < 3:
        tx(sl, lf+sw+0.02, 1.65, 0.13, 0.55, '>', 24, TEAL, bold=True, al=PP_ALIGN.CENTER)
benefits = [
    ('⏱', '對帳時間縮短 50%+', '月底對帳從數小時縮至 30 分鐘以內'),
    ('🎯', '準確率大幅提升', 'AI 不會視覺疲勞，每筆逐一核對'),
    ('📊', '自動輸出 Excel', '格式統一，立即可用，方便存檔'),
    ('🔍', '快速找出差異', '金額、數量異常自動標記，一目了然'),
]
for i, (icon, title, detail) in enumerate(benefits):
    col = i % 2; row = i // 2
    lf = 0.35 + col * 6.5; t = 3.4 + row * 1.6
    bg = OFFWH if i % 4 in [0, 3] else RGBColor(0xE8, 0xF5, 0xF8)
    rc(sl, lf, t, 6.2, 1.45, bg)
    rc(sl, lf, t, 0.8, 1.45, TEAL if col == 0 else SEAFOAM)
    tx(sl, lf+0.1, t+0.35, 0.62, 0.72, icon, 24, WHITE, al=PP_ALIGN.CENTER)
    tx(sl, lf+0.95, t+0.08, 5.1, 0.48, title, 16, NAVY, bold=True)
    tx(sl, lf+0.95, t+0.55, 5.1, 0.60, detail, 13, MIDGRAY)

# ─── SLIDE 5: Results ────────────────────────────────────────────────────────
sl = bl()
hdr(sl, '本週成果 — 缺料分析系統', 4)
sw2 = (W-0.5)/3 - 0.13
for i, (num, lbl, clr) in enumerate([
    ('150+', '項原物料追蹤', TEAL),
    ('30秒', '完成全套分析', SEAFOAM),
    ('3 檔', '資料自動整合', NAVY),
]):
    lf = 0.25 + i * (sw2 + 0.13)
    rc(sl, lf, 1.1, sw2, 1.6, clr)
    tc = WHITE if clr != SEAFOAM else DKGRAY
    tx(sl, lf, 1.17, sw2, 0.95, num, 44, tc, bold=True, al=PP_ALIGN.CENTER)
    tx(sl, lf, 2.02, sw2, 0.50, lbl, 14, tc, al=PP_ALIGN.CENTER)
rc(sl, 0.25, 2.9, W-0.45, 0.42, NAVY)
tx(sl, 0.42, 2.97, W-0.65, 0.30, '▌ 系統功能說明', 14, WHITE, bold=True)
func_items = [
    ('結合訂單總表 × 原物料庫存 × BOM 三檔資料自動整合'),
    ('自動計算每筆訂單缺料項目與缺少數量'),
    ('輸出缺料清單，標示急迫程度（嚴重 / 高 / 中）'),
    ('已追蹤布料、耳繩、印花外層 150+ 項原物料'),
]
fi = ['📑', '⚙', '🚦', '🧵']
for i, text in enumerate(func_items):
    t = 3.45 + i * 0.72
    bg = OFFWH if i % 2 == 0 else WHITE
    rc(sl, 0.25, t, W-0.45, 0.65, bg)
    tx(sl, 0.36, t+0.08, 0.52, 0.52, fi[i], 20, TEAL, al=PP_ALIGN.CENTER)
    tx(sl, 0.98, t+0.10, W-1.2, 0.45, text, 15, DKGRAY)
rc(sl, 0.25, 6.4, W-0.45, 0.90, TEAL)
tx(sl, 0.45, 6.47, 5.5, 0.70, '手動分析：約 2 小時', 18, WHITE, bold=True)
tx(sl, 6.8, 6.47, 1.2, 0.70, '→', 28, GOLD, bold=True, al=PP_ALIGN.CENTER)
tx(sl, 8.2, 6.47, 4.8, 0.70, 'AI 分析：30 秒以內', 18, WHITE, bold=True)

# ─── SLIDE 6: Insights ───────────────────────────────────────────────────────
sl = bl()
hdr(sl, '使用 Claude 的心得與體會', 5)
insight_data = [
    ('💡', 'AI 不是萬能，需要清楚描述需求',
     '說得越清楚，AI 做得越準確；業務邏輯說明清楚是關鍵'),
    ('📈', '從「告訴 AI 做什麼」進步到「讓 AI 理解業務邏輯」',
     '從一步步指令，到讓 AI 自主判斷 — 學習曲線明顯'),
    ('🧠', '永久記憶功能讓每次對話延續',
     '不用重複說明，AI 記得採購流程與格式規範'),
    ('🚀', '最大收穫：讓重複性工作自動化',
     '核對、整理、彙整交給 AI，人專注在判斷與決策'),
]
acs = [TEAL, SEAFOAM, NAVY, TEALBG]
for i, (icon, title, detail) in enumerate(insight_data):
    t = 1.15 + i * 1.35
    rc(sl, 0.3, t, 0.18, 1.20, acs[i])
    rc(sl, 0.48, t, W-0.70, 1.20, OFFWH if i % 2 == 0 else WHITE)
    rc(sl, 0.65, t+0.22, 0.72, 0.72, acs[i])
    tx(sl, 0.65, t+0.22, 0.72, 0.72, icon, 28, WHITE, al=PP_ALIGN.CENTER)
    tx(sl, 1.55, t+0.08, W-1.90, 0.45, title, 15, NAVY, bold=True)
    tx(sl, 1.55, t+0.55, W-1.90, 0.52, detail, 13, MIDGRAY)

# ─── SLIDE 7: Future Plans ───────────────────────────────────────────────────
sl = bl()
hdr(sl, '下一步 — 繼續擴大應用', 6)
plan_data = [
    ('01', '採購比價自動化',        '多廠商報價一鍵比較，自動產出推薦清單',       TEAL),
    ('02', '進貨入庫通知整合',      '入庫資料即時核對，異常自動提醒',             NAVY),
    ('03', '與 ERP 系統深度串接',   '直接讀取鼎新資料，減少手動匯出步驟',        SEAFOAM),
    ('04', '建立標準化 AI SOP',     '讓全團隊都能上手，複製成功經驗',            TEALBG),
]
for i, (num, title, detail, clr) in enumerate(plan_data):
    col = i % 2; row = i // 2
    lf = 0.35 + col * 6.5; t = 1.2 + row * 2.7
    rc(sl, lf, t, 6.2, 2.4, clr)
    tx(sl, lf+0.15, t+0.08, 1.5, 0.90, num, 42, GOLD, bold=True)
    tx(sl, lf+0.15, t+1.0, 5.8, 0.58, title, 18, WHITE, bold=True)
    tc = WHITE if clr != SEAFOAM else DKGRAY
    tx(sl, lf+0.15, t+1.60, 5.8, 0.65, detail, 13, LTBLUE if clr != SEAFOAM else DKGRAY)

# ─── SLIDE 8: Closing ────────────────────────────────────────────────────────
sl = bl()
rc(sl, 0, 0, W, H, DNAV)
rc(sl, 0, 0, 0.35, H, TEAL)
rc(sl, 0, H-1.0, W, 1.0, TEALBG)
rc(sl, 0.35, 2.7, W-0.35, 0.05, TEAL)
tx(sl, 0.65, 0.65, W-1.0, 1.05, 'AI 是工具，判斷還是人', 36, WHITE, bold=True)
tx(sl, 0.65, 1.80, W-1.0, 0.70,
   '讓繁瑣的資料整理交給 AI，把時間留給更重要的決策', 18, LTBLUE)
closing = [
    'Claude 幫助處理繁瑣的資料整理工作',
    '讓採購人員有更多時間做供應商管理與談判',
    '歡迎大家分享各自的使用經驗，一起學習進步',
    '一起讓工作更有效率！',
]
ck = ['✅', '✅', '💬', '🚀']
for i, txt in enumerate(closing):
    t = 2.9 + i * 0.82
    tx(sl, 0.80, t, 0.50, 0.65, ck[i], 20, SEAFOAM, al=PP_ALIGN.CENTER)
    tx(sl, 1.45, t+0.08, W-2.0, 0.52, txt, 16, WHITE)
tx(sl, 0.65, H-0.82, 8.0, 0.40,
   '廠務助理／採購  ADA　　達特世　　2026 年 4 月', 13, LTBLUE)

# ─── Save ─────────────────────────────────────────────────────────────────────
out = r'C:\Users\admin\OneDrive\桌面\CODE資料\AI週會心得報告_0422.pptx'
prs.save(out)
print('Saved:', out)
print('Slides:', len(prs.slides))
