# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BG   = RGBColor(0x2C, 0x26, 0x54)
PINK      = RGBColor(0xE8, 0xA0, 0xB4)
PURPLE    = RGBColor(0x9B, 0x8E, 0xC4)
LIGHT_BG  = RGBColor(0xF5, 0xF3, 0xFA)
RED_SEAL  = RGBColor(0xC0, 0x39, 0x2B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY      = RGBColor(0xCC, 0xCC, 0xCC)
LIGHT_PUR = RGBColor(0xE8, 0xE6, 0xF5)
LIGHT_PIK = RGBColor(0xF9, 0xE8, 0xEF)
LIGHT_GRN = RGBColor(0xEB, 0xF7, 0xEE)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
DARK_GRN  = RGBColor(0x2C, 0x7A, 0x3E)
LIGHT_BLU = RGBColor(0xE8, 0xF4, 0xF8)
BLUE      = RGBColor(0x5D, 0xAD, 0xE2)
DARK_BLU  = RGBColor(0x1A, 0x5C, 0x8A)
AMBER     = RGBColor(0xE6, 0x7E, 0x22)
TEAL      = RGBColor(0x17, 0xA5, 0x8A)
LIGHT_TEL = RGBColor(0xE0, 0xF5, 0xF1)

# ── 16:9 尺寸 ──
W = 13.33
H = 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def box(s, x, y, w, h, fc=None, bc=None, lw=1.5):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else: sh.fill.background()
    if bc: sh.line.color.rgb = bc; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def circ(s, x, y, d, fc=None):
    sh = s.shapes.add_shape(9, Inches(x), Inches(y), Inches(d), Inches(d))
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else: sh.fill.background()
    sh.line.fill.background()
    return sh

def txt(s, t, x, y, w, h, sz, c, bold=False, italic=False, align=PP_ALIGN.LEFT):
    """文字框 - x/y/w/h 皆為 inch，sz 為 pt"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.line.fill.background()          # ← 移除文字框邊框
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = t
    r.font.size = Pt(sz)
    r.font.color.rgb = c
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Microsoft JhengHei"
    return tb

# 文字垂直置中輔助：在色塊 (bx,by,bw,bh) 內，依字體大小自動計算 y
def ctxt(s, t, bx, by, bw, bh, sz, c, bold=False, italic=False, align=PP_ALIGN.CENTER,
         pad_x=0.15, lines=1):
    """色塊內置中文字：pad_x=左右內距, lines=預估行數"""
    line_h = sz * 0.016          # 1pt ≈ 0.0139 inch，加緩衝
    th = line_h * lines + 0.05   # 文字框高度
    ty = by + (bh - th) / 2      # 垂直置中
    tx = bx + pad_x
    tw = bw - pad_x * 2
    return txt(s, t, tx, ty, tw, th, sz, c, bold=bold, italic=italic, align=align)

def deco(s):
    box(s, W-3.0, 0, 3.0, 1.0, fc=PINK)

def left_bar(s, c):
    box(s, 0, 0, 0.15, H, fc=c)

def slide_title(s, num, title, bar_c):
    left_bar(s, bar_c)
    deco(s)
    box(s, 0.35, 0.18, 0.65, 0.65, fc=bar_c)
    ctxt(s, num, 0.35, 0.18, 0.65, 0.65, 20, WHITE, bold=True, lines=1)
    txt(s, title, 1.2, 0.22, 11.5, 0.58, 26, DARK_BG, bold=True)

# ══════════════════════════════
# 1. 封面
# ══════════════════════════════
s = blank()
bg(s, DARK_BG)
deco(s)
box(s, 0, 3.8, W, 0.08, fc=PINK)
box(s, 0, 3.92, W, 0.04, fc=PURPLE)
box(s, 0.5, H-0.5, 0.45, 0.38, fc=RED_SEAL)

txt(s, "AI 應用心得週報", 0.5, 1.2, W-1, 1.1, 44, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "兩週完整整合報告  |  2026 年 4 月", 0.5, 2.55, W-1, 0.65, 18, PINK, align=PP_ALIGN.CENTER)

topics = [
    (PINK,   "採購對帳 AI 實驗"),
    (PURPLE, "資料夾整理"),
    (BLUE,   "比價換算單位"),
    (TEAL,   "AI 使用量分析"),
]
for i, (c, label) in enumerate(topics):
    x = 1.5 + i * 2.6
    box(s, x, 4.15, 2.3, 0.75, fc=c)
    ctxt(s, label, x, 4.15, 2.3, 0.75, 16, WHITE, bold=True, lines=1)

txt(s, "廠務助理／採購　ADA", 0.5, 5.45, W-1, 0.55, 16, WHITE, align=PP_ALIGN.CENTER)
txt(s, "達特世生技有限公司", 0.5, 6.08, W-1, 0.5, 14, GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════
# 2. 目錄
# ══════════════════════════════
s = blank()
bg(s, LIGHT_BG)
deco(s)
txt(s, "本次報告內容", 0.7, 0.22, 9, 0.7, 28, DARK_BG, bold=True)
txt(s, "兩週 AI 使用心得整合", 0.7, 0.85, 9, 0.42, 15, PURPLE, italic=True)

items = [
    (PINK,   "1", "上週回顧：採購對帳 AI 實驗",  "截圖→AI比對，找出喬煒少入庫 46 Kg"),
    (PURPLE, "2", "本週一：資料夾整理",           "重複檔案刪除，節省電腦儲存空間"),
    (BLUE,   "3", "本週二：比價換算單位計價",     "不同單位報價快速換算比較（昊瑞 vs 全程）"),
    (GREEN,  "4", "本週三：AI 協助新品編碼",        "提供裝箱單→AI依編碼規則生成7筆新品號"),
    (TEAL,   "5", "本週四：AI 使用量自我分析",    "盤點各項 AI 使用，缺料分析已成每日例行"),
    (DARK_BG,"6", "整體學習心得與下一步",         "AI 最大功用：創造工具，不只是問答"),
]
for i, (c, num, title, sub) in enumerate(items):
    y = 1.30 + i * 0.94   # 6 items × 0.94 = 5.64", ends at 7.00" ✓
    bh = 0.88
    box(s, 0.45, y, W-0.9, bh, fc=WHITE, bc=c)
    circ(s, 0.6, y + bh/2 - 0.28, 0.56, fc=c)
    ctxt(s, num, 0.6, y + bh/2 - 0.28, 0.56, 0.56, 16, WHITE, bold=True, lines=1)
    txt(s, title, 1.42, y+0.08, W-2.2, 0.38, 16, DARK_BG, bold=True)
    txt(s, sub,   1.42, y+0.48, W-2.2, 0.32, 13, DARK_TEXT, bold=True)

# ══════════════════════════════
# 3. 上週回顧（採購對帳）
# ══════════════════════════════
s = blank()
bg(s, LIGHT_BG)
slide_title(s, "1", "上週回顧：採購對帳 AI 實驗", PINK)

bh = 0.58
box(s, 0.35, 1.05, W-0.7, bh, fc=DARK_BG)
ctxt(s, "人工對帳 → 沒發現異常  →  改用 AI  →  找出少入庫 46 Kg！",
     0.35, 1.05, W-0.7, bh, 15, PINK, bold=True, lines=1)

# 兩大實驗
exp_data = [
    (PURPLE, LIGHT_PUR, "實驗一：鼎新截圖",
     "操作：鼎新系統預覽列印截圖\n→ 存入資料夾 → 告知 AI 資料夾路徑\n→ AI 自行讀取並比對\n\n對象：喬煒興業（耳繩採購）\n結果：發現少入庫 46 Kg！\n       差額含稅 10,143 元"),
    (PINK, LIGHT_PIK, "實驗二：拍照 → LINE → AI",
     "操作：手機拍照採購單/發票\n→ 用 LINE 上傳到資料夾\n→ AI 自行讀取資料夾\n\n對象：凱達印刷（包材採購）\n結果：解析度足夠！\n       採購單與進貨單全部相符"),
]
EW = (W - 0.35*2 - 0.2) / 2   # 每欄寬度
for i, (hc, bgc, title, desc) in enumerate(exp_data):
    x = 0.35 + i * (EW + 0.2)
    box(s, x, 1.82, EW, 0.52, fc=hc)
    ctxt(s, title, x, 1.82, EW, 0.52, 15, WHITE, bold=True, lines=1)
    box(s, x, 2.34, EW, 2.95, fc=bgc)
    txt(s, desc, x+0.18, 2.44, EW-0.36, 2.75, 14, DARK_TEXT, bold=True)

# 成果數字
stats = [
    (RED_SEAL, "46 Kg",  "少入庫量"),
    (AMBER,    "10,143", "含稅差額(元)"),
    (GREEN,    "30 秒",  "AI完成時間"),
    (DARK_BG,  "2 → 0",  "對帳差錯次數"),
]
SW = (W - 0.35*2 - 0.15*3) / 4
for i, (c, num, label) in enumerate(stats):
    x = 0.35 + i * (SW + 0.15)
    box(s, x, 5.48, SW, 1.7, fc=c)
    txt(s, num,   x+0.1, 5.62, SW-0.2, 0.8, 28, WHITE, bold=True, align=PP_ALIGN.CENTER)
    ctxt(s, label, x, 6.45, SW, 0.6, 14, WHITE, bold=True, lines=1)

# ══════════════════════════════
# 3b. 喬煒差異單據詳情
# ══════════════════════════════
s = blank()
bg(s, LIGHT_BG)
left_bar(s, RED_SEAL)
deco(s)
txt(s, "實驗一成果：喬煒興業  截圖 → AI 找出差異", 0.55, 0.18, 12.2, 0.62, 26, DARK_BG, bold=True)

# 橫幅
bh = 0.56
box(s, 0.35, 0.98, W-0.7, bh, fc=DARK_BG)
ctxt(s, "鼎新系統截圖 → AI 自行讀取資料夾 → 人工對帳未發現，AI 找出：少入庫 46 Kg！",
     0.35, 0.98, W-0.7, bh, 15, PINK, bold=True, lines=1)

# 欄寬
CW2 = (W - 0.35*2 - 0.25) / 2

# 採購單（左）
box(s, 0.35, 1.72, CW2, 0.5, fc=PURPLE)
ctxt(s, "採購單　3301-20260223002", 0.35, 1.72, CW2, 0.5, 15, WHITE, bold=True,
     align=PP_ALIGN.LEFT, pad_x=0.2, lines=1)
box(s, 0.35, 2.22, CW2, 3.05, fc=LIGHT_PUR)
po_items = [
    ("品號",   "BEL-JW037-WHT500",  False),
    ("品名",   "耳繩-喬煒3.7-白500", False),
    ("廠商",   "喬煒興業有限公司",    False),
    ("單據日期","2026/02/23",         False),
    ("採購數量","178 Kg",             True),
    ("採購單價","210 元/Kg",          False),
    ("採購金額","37,380 元",          True),
    ("含稅合計","39,249 元",          True),
]
for i, (k, v, hi) in enumerate(po_items):
    y2 = 2.32 + i * 0.36
    txt(s, k + "：", 0.52, y2, 2.0, 0.34, 14, DARK_TEXT, bold=True)
    txt(s, v, 2.7, y2, CW2-2.5, 0.34, 14, PURPLE if hi else DARK_TEXT, bold=hi)

# 進貨單（右）
RX2 = 0.35 + CW2 + 0.25
box(s, RX2, 1.72, CW2, 0.5, fc=RED_SEAL)
ctxt(s, "進貨單　20260226002", RX2, 1.72, CW2, 0.5, 15, WHITE, bold=True,
     align=PP_ALIGN.LEFT, pad_x=0.2, lines=1)
box(s, RX2, 2.22, CW2, 3.05, fc=LIGHT_PIK)
gr_items = [
    ("品號",    "BEL-JW037-WHT500",    False),
    ("品名",    "耳繩-喬煒3.7-白500",   False),
    ("供應商",  "喬煒興業有限公司",      False),
    ("單據日期","2026/02/26",            False),
    ("實收數量","132 Kg  ⚠",            True),
    ("計價單價","210 元/Kg",             False),
    ("進貨金額","27,720 元  ⚠",         True),
    ("含稅合計","29,106 元  ⚠",         True),
]
for i, (k, v, hi) in enumerate(gr_items):
    y2 = 2.32 + i * 0.36
    txt(s, k + "：", RX2+0.18, y2, 2.0, 0.34, 14, DARK_TEXT, bold=True)
    txt(s, v, RX2+2.35, y2, CW2-2.5, 0.34, 14, RED_SEAL if hi else DARK_TEXT, bold=hi)

# 差異明細
YELLOW = RGBColor(0xFF, 0xF3, 0xCD)
box(s, 0.35, 5.42, W-0.7, 0.48, fc=AMBER)
ctxt(s, "⚠  差異明細", 0.35, 5.42, W-0.7, 0.48, 15, WHITE, bold=True, lines=1)
box(s, 0.35, 5.9, W-0.7, 1.0, fc=YELLOW)
diff_items = [
    ("數量差異：",      "採購 178 Kg  →  實收 132 Kg",  "少 46 Kg"),
    ("金額差異（未稅）：","37,380 元  →  27,720 元",      "差 9,660 元"),
    ("金額差異（含稅）：","39,249 元  →  29,106 元",      "差 10,143 元"),
]
for i, (label, detail, diff) in enumerate(diff_items):
    y2 = 5.98 + i * 0.3
    txt(s, label,  0.52, y2, 3.2, 0.28, 13, DARK_TEXT, bold=True)
    txt(s, detail, 3.8,  y2, 5.8, 0.28, 13, DARK_TEXT, bold=True)
    txt(s, diff,   9.75, y2, 2.9, 0.28, 14, RED_SEAL,  bold=True, align=PP_ALIGN.RIGHT)

box(s, 0.35, 7.0, W-0.7, 0.38, fc=DARK_BG)
ctxt(s, "後續更正：補開一張採購進貨單，補入 46 Kg（金額 9,660 元，含稅 10,143 元）",
     0.35, 7.0, W-0.7, 0.38, 14, GREEN, bold=True, lines=1)

# ══════════════════════════════
# 4. 本週一：資料夾整理
# ══════════════════════════════
s = blank()
bg(s, LIGHT_BG)
slide_title(s, "2", "本週一：資料夾整理（重複檔案刪除）", PURPLE)

bh = 0.56
box(s, 0.35, 1.05, W-0.7, bh, fc=PURPLE)
ctxt(s, "問題：電腦存了很多重複或類似的檔案，佔用大量儲存空間",
     0.35, 1.05, W-0.7, bh, 15, WHITE, bold=True, lines=1)

# 左右欄
CW = (W - 0.35*2 - 0.2) / 2

# 左：傳統困擾
hh = 0.5
box(s, 0.35, 1.78, CW, hh, fc=RED_SEAL)
ctxt(s, "以前的困擾", 0.35, 1.78, CW, hh, 15, WHITE, bold=True, lines=1)
prob_list = [
    "× 同一份報表存了很多版本",
    "× 不確定哪個才是最新版",
    "× 怕刪錯不敢輕易刪除",
    "× 資料夾越來越亂、越來越大",
    "× 搜尋檔案要花很多時間",
]
box(s, 0.35, 2.28, CW, 3.0, fc=LIGHT_PIK)
for i, p in enumerate(prob_list):
    txt(s, p, 0.52, 2.40+i*0.55, CW-0.3, 0.48, 14, DARK_TEXT, bold=True)

# 右：AI協助方式
RX = 0.35 + CW + 0.2
box(s, RX, 1.78, CW, hh, fc=DARK_GRN)
ctxt(s, "AI 協助方式", RX, 1.78, CW, hh, 15, WHITE, bold=True, lines=1)
ai_list = [
    "√ 說明資料夾結構，讓 AI 分析",
    "√ AI 列出可能重複的檔案清單",
    "√ AI 協助判斷哪些版本可刪除",
    "√ 建議清楚的命名規則（含日期）",
    "√ 整理後資料夾層次更清晰",
]
box(s, RX, 2.28, CW, 3.0, fc=LIGHT_GRN)
for i, a in enumerate(ai_list):
    txt(s, a, RX+0.18, 2.40+i*0.55, CW-0.3, 0.48, 14, DARK_TEXT, bold=True)

# 成果 3 格
box(s, 0.35, 5.45, W-0.7, 0.46, fc=DARK_BG)
ctxt(s, "本週成果", 0.35, 5.45, W-0.7, 0.46, 14, PINK, bold=True, align=PP_ALIGN.LEFT,
     pad_x=0.2, lines=1)

RW3 = (W - 0.35*2 - 0.15*2) / 3
results = [
    (LIGHT_PUR, PURPLE,  "整理方式",  "請 AI 協助規劃資料夾\n命名與分類原則"),
    (LIGHT_GRN, DARK_GRN,"節省空間",  "減少重複檔案佔用的\n硬碟儲存容量"),
    (LIGHT_BLU, DARK_BLU,"命名規則",  "統一「YYYYMM_品名_版本」\n的格式，易於搜尋"),
]
for i, (bgc, tc, title, desc) in enumerate(results):
    x = 0.35 + i * (RW3 + 0.15)
    box(s, x, 6.0, RW3, 1.35, fc=bgc)
    txt(s, title, x+0.15, 6.08, RW3-0.3, 0.4, 14, tc, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  x+0.15, 6.52, RW3-0.3, 0.75, 13, DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════
# 5. 本週二：比價換算
# ══════════════════════════════
s = blank()
bg(s, LIGHT_BG)
slide_title(s, "3", "本週二：協助比價換算單位計價", BLUE)

bh = 0.56
box(s, 0.35, 1.05, W-0.7, bh, fc=DARK_BLU)
ctxt(s, "問題：兩家廠商報價單位不同（USD/Kg vs NT$/m²），無法直接比較，需要 AI 協助換算",
     0.35, 1.05, W-0.7, bh, 14, WHITE, bold=True, lines=1)

box(s, 0.35, 1.72, W-0.7, 0.48, fc=BLUE)
ctxt(s, "實際案例：採購彈性布料，昊瑞（USD計價） vs 全程（NT$/m²計價），請 AI 換算後比較",
     0.35, 1.72, W-0.7, 0.48, 14, WHITE, lines=1)

# 表格（2 家廠商）
HDR_H = 0.5
ROW_H = 0.9
hdrs = ["廠商", "原始報價", "換算條件", "換算後（NT$/Kg）", "結論"]
ws = [1.6, 2.5, 3.5, 2.6, 1.08]
xs = [0.35]
for w2 in ws[:-1]:
    xs.append(xs[-1] + w2 + 0.05)

box(s, 0.35, 2.32, W-0.7, HDR_H, fc=DARK_BG)
for j, (h, w2, x2) in enumerate(zip(hdrs, ws, xs)):
    ctxt(s, h, x2, 2.32, w2, HDR_H, 14, WHITE, bold=True, lines=1)

rows = [
    ("昊瑞",  "USD $4.18/Kg", "× 1.2運費 × 33匯率",               "NT$ 165.5/Kg", "較低"),
    ("全程",  "NT$18/m²",    "100,000m²=5,800Kg → 1Kg≈17.24m²", "NT$ 310.3/Kg", "較高"),
]
row_bgs = [LIGHT_GRN, LIGHT_BLU]
row_tcs = [DARK_GRN,  DARK_BLU]
for ri, (row, rbg, rtc) in enumerate(zip(rows, row_bgs, row_tcs)):
    y2 = 2.32 + HDR_H + ri * ROW_H + 0.04
    box(s, 0.35, y2, W-0.7, ROW_H-0.04, fc=rbg)
    for j, (cell, w2, x2) in enumerate(zip(row, ws, xs)):
        c2 = rtc if j == 3 else DARK_TEXT
        bold2 = (j == 3)
        sz2 = 14 if j != 2 else 13
        ctxt(s, cell, x2, y2, w2, ROW_H-0.04, sz2, c2, bold=bold2, lines=2 if j==2 else 1)

# AI 步驟
SY = 2.32 + HDR_H + 2*ROW_H + 0.15
box(s, 0.35, SY, W-0.7, 0.48, fc=PURPLE)
ctxt(s, "AI 協助步驟", 0.35, SY, W-0.7, 0.48, 15, WHITE, bold=True,
     align=PP_ALIGN.LEFT, pad_x=0.2, lines=1)

steps = [
    (LIGHT_PUR, "提供兩家廠商\n的報價資料"),
    (LIGHT_BLU, "告訴 AI 換算\n條件（匯率/運費）"),
    (LIGHT_GRN, "AI 自動換算\n到相同單位"),
    (LIGHT_PIK, "AI 給出比較\n結論與建議"),
]
SW2 = (W - 0.35*2 - 0.15*3) / 4
for i, (bgc, desc) in enumerate(steps):
    x2 = 0.35 + i * (SW2 + 0.15)
    sy2 = SY + 0.55
    box(s, x2, sy2, SW2, 1.1, fc=bgc)
    txt(s, desc, x2+0.15, sy2+0.15, SW2-0.3, 0.8, 14, DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, "→", x2+SW2+0.02, sy2+0.3, 0.13, 0.45, 18, PURPLE, bold=True)

EY = SY + 0.55 + 1.1 + 0.1
box(s, 0.35, EY, W-0.7, 0.88, fc=LIGHT_BLU)
txt(s, "效益：", 0.52, EY+0.1, 1.4, 0.38, 14, DARK_BLU, bold=True)
txt(s, "跨幣別（USD→NT$）+ 跨單位（m²→Kg）計算，以前需要多個計算機來回算\n只要把廠商報價和換算條件告訴 AI，幾秒鐘就能完成，不用怕算錯！",
    1.85, EY+0.1, W-2.3, 0.72, 13, DARK_TEXT, bold=True)

# ══════════════════════════════
# 5b. 本週三：AI 協助新品編碼
# ══════════════════════════════
s = blank()
bg(s, LIGHT_BG)
slide_title(s, "4", "本週三：AI 協助新品編碼 — 綠沅 6.0mm 耳繩", GREEN)

bh = 0.56
box(s, 0.35, 1.05, W-0.7, bh, fc=DARK_GRN)
ctxt(s, "提供裝箱單（PDF）給 AI → AI 讀取顏色與色碼 → 依編碼原則自動生成 7 筆新品號",
     0.35, 1.05, W-0.7, bh, 15, WHITE, bold=True, lines=1)

# 編碼格式說明（色塊加高，確保標注文字不溢出）
BEL_BOX_Y = 1.72
BEL_BOX_H = 0.75   # 加高：色塊結束於 2.47，標注文字結束於 2.40
box(s, 0.35, BEL_BOX_Y, W-0.7, BEL_BOX_H, fc=LIGHT_GRN)
txt(s, "編碼格式：", 0.55, 1.78, 2.0, 0.32, 14, DARK_GRN, bold=True)
txt(s, "BEL  -  LY  060  -  [顏色3碼][廠商色碼3碼]", 2.4, 1.78, 9.5, 0.32, 15, DARK_BG, bold=True)
txt(s, "耳繩      綠沅  6.0mm       顏色代號   廠商色碼", 2.4, 2.13, 9.5, 0.25, 11, DARK_TEXT)

# 表格標題（下移至色塊結束後）
HDR_Y2 = 2.55
HDR_H = 0.45
box(s, 0.35, HDR_Y2, W-0.7, HDR_H, fc=DARK_BG)
col_hdrs = ["廠商色碼", "中文名稱", "英文描述（裝箱單）", "顏色代號", "品號（新）", "狀態"]
col_ws   = [1.4, 1.5, 3.8, 1.5, 2.8, 1.5]
col_xs   = [0.35]
for w2 in col_ws[:-1]:
    col_xs.append(col_xs[-1] + w2 + 0.04)
for h2, w2, x2 in zip(col_hdrs, col_ws, col_xs):
    ctxt(s, h2, x2, HDR_Y2, w2, HDR_H, 13, WHITE, bold=True, lines=1)

# 資料列
code_rows = [
    ("L342", "新粉紅", "L342 Pink Baiyuan Earloop",       "PNK", "BEL-LY060-PNK342", "新"),
    ("M064", "灰",    "M064 Grey Baiyuan Earloop",        "GRY", "BEL-LY060-GRY064", "新"),
    ("M190", "淺灰",  "M190 Light Grey Blue Baiyuan Earloop","GRY","BEL-LY060-GRY190","新"),
    ("L262", "粉紫",  "L262 Pink Purple Baiyuan Earloop", "PUR", "BEL-LY060-PUR262", "新"),
    ("D335", "鐵灰",  "D335 Iron Grey Baiyuan Earloop",   "GRY", "BEL-LY060-GRY335", "新"),
    ("L236", "冷灰藍","L236 Cool Grey Blue Baiyuan Earloop","BLU","BEL-LY060-BLU236","新"),
    ("L157", "淺紫",  "L157 Light Purple Baiyuan Earloop","PUR", "BEL-LY060-PUR157", "新"),
]
ROW_H2 = 0.50
row_bgs2 = [LIGHT_PIK, WHITE, LIGHT_PIK, WHITE, LIGHT_PIK, WHITE, LIGHT_PIK]
for ri, (row, rbg) in enumerate(zip(code_rows, row_bgs2)):
    y2 = HDR_Y2 + HDR_H + ri * ROW_H2 + 0.04
    box(s, 0.35, y2, W-0.7, ROW_H2-0.04, fc=rbg)
    for j, (cell, w2, x2) in enumerate(zip(row, col_ws, col_xs)):
        c2 = DARK_GRN if j == 4 else (GREEN if j == 5 else DARK_TEXT)
        bold2 = j in [4, 5]
        ctxt(s, cell, x2, y2, w2, ROW_H2-0.04, 13 if j != 2 else 11, c2, bold=bold2, lines=1 if j != 2 else 2)

# 底部說明
box(s, 0.35, H-0.72, W-0.7, 0.58, fc=LIGHT_GRN)
txt(s, "已確認：", 0.55, H-0.64, 1.5, 0.35, 14, DARK_GRN, bold=True)
txt(s, "灰色系（064 / 190 / 335）使用 GRY，藍灰系（236）使用 BLU，紫色系（157 / 262）使用 PUR，粉色（342）使用 PNK\n明日輸入 ERP 建立新品號，完成建檔",
    1.9, H-0.64, W-2.3, 0.5, 13, DARK_TEXT, bold=True)

# ══════════════════════════════
# 6. 本週三：AI 使用量分析
# ══════════════════════════════
s = blank()
bg(s, LIGHT_BG)
slide_title(s, "5", "本週四：AI 使用量自我分析", TEAL)

bh = 0.56
box(s, 0.35, 1.05, W-0.7, bh, fc=TEAL)
ctxt(s, "盤點這兩週 AI 的使用情況，看看 Claude 實際幫到了哪些工作",
     0.35, 1.05, W-0.7, bh, 15, WHITE, bold=True, lines=1)

box(s, 0.35, 1.73, W-0.7, 0.42, fc=DARK_BG)
ctxt(s, "目前 AI 使用類別分析", 0.35, 1.73, W-0.7, 0.42, 14, PINK, bold=True, lines=1)

uses = [
    (PINK,   "採購對帳",   "截圖/拍照→AI比對\n採購單與進貨單",       "高頻使用"),
    (PURPLE, "PPT 製作",   "把工作內容說給AI\nAI 協助生成投影片",     "高頻使用"),
    (BLUE,   "比價換算",   "不同單位報價換算\n找出最划算廠商",         "中頻使用"),
    (TEAL,   "資料夾整理", "請AI協助規劃命名\n與分類，清除重複",       "中頻使用"),
    (GREEN,  "缺料分析",   "每日執行缺料分析\n新增優先備料排序",       "每日使用"),
    (AMBER,  "問題解答",   "工作問題查詢\n換算說明",                   "日常使用"),
]
UW = (W - 0.35*2 - 0.1*2) / 3
UH = 1.95
for i, (c, title, desc, freq) in enumerate(uses):
    col = i % 3
    row = i // 3
    x2 = 0.35 + col * (UW + 0.1)
    y2 = 2.28 + row * (UH + 0.1)
    box(s, x2, y2, UW, UH, fc=WHITE, bc=c, lw=2)
    box(s, x2, y2, UW, 0.48, fc=c)
    ctxt(s, title, x2, y2, UW, 0.48, 15, WHITE, bold=True, lines=1)
    txt(s, desc, x2+0.15, y2+0.58, UW-0.3, 0.9, 14, DARK_TEXT, bold=True)
    box(s, x2+0.12, y2+UH-0.38, 1.3, 0.3, fc=c)
    ctxt(s, freq, x2+0.12, y2+UH-0.38, 1.3, 0.3, 12, WHITE, bold=True, lines=1)

box(s, 0.35, H-0.88, W-0.7, 0.75, fc=DARK_BG)
txt(s, "使用趨勢：", 0.55, H-0.82, 1.8, 0.35, 14, PINK, bold=True)
txt(s, "從初期只會問答，到現在能用 AI 執行採購對帳、缺料分析、比價換算等實際工作任務，每個工作日都有在用，使用範圍持續擴大中",
    2.2, H-0.82, W-2.6, 0.62, 13, WHITE, bold=True)

# ══════════════════════════════
# 7. 整體學習心得
# ══════════════════════════════
s = blank()
bg(s, LIGHT_BG)
slide_title(s, "5", "整體學習心得與下一步", PINK)

insights = [
    (PURPLE, "指令說清楚，結果就越準確",
             "告訴 AI「哪個廠商、哪些檔案、什麼規則、要什麼格式」\n說越清楚，AI 給的結果越精準，來回修正的次數越少"),
    (BLUE,   "AI 擔任二次確認角色  →  降低對帳錯誤率",
             "以往痛點：對帳完沒有人再次確認單據，錯誤容易被忽略\n實際案例：人工對完帳沒發現異常，AI 二次檢查抓出少入庫 46 Kg\n→ 有 AI 幫忙把關，單據錯誤率明顯降低"),
    (TEAL,   "下一步：繼續擴大應用",
             "每月月初對帳流程 → 全部用 AI 輔助完成\n新廠商報價比價 → 標準化成固定的 AI 協助流程\n缺料分析 → 持續優化優先備料排序功能"),
]
IH = 1.65
for i, (c, title, desc) in enumerate(insights):
    y2 = 1.05 + i * (IH + 0.12)
    bg_map = {PURPLE: LIGHT_PUR, BLUE: LIGHT_BLU, TEAL: LIGHT_TEL}
    bgc = bg_map.get(c, LIGHT_PIK)
    box(s, 0.35, y2, 0.14, IH, fc=c)
    box(s, 0.52, y2, W-0.87, IH, fc=bgc)
    txt(s, title, 0.72, y2+0.12, W-1.2, 0.48, 16, DARK_BG, bold=True)
    txt(s, desc,  0.72, y2+0.65, W-1.2, 0.9, 14, DARK_TEXT, bold=True)

# ══════════════════════════════
# 8. 結語
# ══════════════════════════════
s = blank()
bg(s, DARK_BG)
deco(s)
box(s, 0, H/2-0.05, W, 0.08, fc=PINK)
box(s, 0.5, H-0.52, 0.45, 0.38, fc=RED_SEAL)

txt(s, "感謝聆聽", 0.5, 1.2, W-1, 1.1, 44, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "AI 是工具，讓我把時間花在更重要的判斷上", 0.5, 2.55, W-1, 0.65, 18, PINK, align=PP_ALIGN.CENTER)

ach = [
    (PINK,   "採購對帳",  "AI 找出少入庫 46 Kg"),
    (PURPLE, "資料夾整理","重複檔案清除節省空間"),
    (BLUE,   "比價換算",  "秒算兩家廠商最划算報價"),
    (TEAL,   "缺料分析",  "新增優先備料排序功能"),
]
AW = (W - 0.5*2 - 0.15*3) / 4
for i, (c, title, desc) in enumerate(ach):
    x2 = 0.5 + i * (AW + 0.15)
    box(s, x2, 3.75, AW, 1.85, fc=RGBColor(0x3D,0x35,0x6B))
    box(s, x2, 3.75, AW, 0.52, fc=c)
    ctxt(s, title, x2, 3.75, AW, 0.52, 16, WHITE, bold=True, lines=1)
    txt(s, desc, x2+0.15, 4.38, AW-0.3, 1.05, 15, GRAY, bold=True, align=PP_ALIGN.CENTER)

txt(s, "廠務助理／採購　ADA", 0.5, 5.95, W-1, 0.55, 16, WHITE, align=PP_ALIGN.CENTER)
txt(s, "2026 年 4 月  |  達特世生技有限公司", 0.5, 6.55, W-1, 0.5, 14, GRAY, align=PP_ALIGN.CENTER)

# 儲存
out = r"C:\Users\admin\OneDrive\桌面\CODE資料\ADA_AI週報_整合版_0429.pptx"
prs.save(out)
print("done:", len(prs.slides), "slides ->", out)
